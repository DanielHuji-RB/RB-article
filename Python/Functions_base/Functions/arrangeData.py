# Daniel Sand

import numpy as np
from sklearn.model_selection import train_test_split



def arrageData_orUseFixData( df,metaDataCol_N,useNewData_Flag ):
    import pickle

    if useNewData_Flag:
        norm_by_subj_Flag=True
        removeOutlaires_Flag=True
        useElec_toSplitGroups_Flag=False
        df_temp=df.iloc[:,metaDataCol_N:].copy()

        if removeOutlaires_Flag:
            df_zscore = (df_temp.iloc[:, 0:] - df_temp.iloc[:, 0:].mean()) / df_temp.iloc[:, 0:].std()
            df_temp[df_zscore.abs() > 2.5] = np.nan
            df_temp_plot = df_temp.copy()


        if useElec_toSplitGroups_Flag:
            df_groups= df_temp.iloc[:,0:].groupby([df['reference'],df['side'],df['Elec']])
        else:
            df_groups = df_temp.iloc[:, 0:].groupby([df['reference'], df['side']])


        if not norm_by_subj_Flag:
            normalized_df=(df_temp-df_temp.mean())/df_temp.std() #normlize by zcocre

        if norm_by_subj_Flag:
            normalized_df=df_groups.apply(lambda x: (x - x.mean())/ (x.std()))


        df.iloc[:,metaDataCol_N:]=normalized_df
        df = df.replace(np.nan, 0)

        if 'TrainOrTest' in df:
            df_train=df.loc[df.TrainOrTest=='Train',:]
            df_test=df.loc[df.TrainOrTest=='Test',:]
        else:
            df_train ,df_test = train_test_split(df,test_size=0.2)

        with open('dataSet_zscore_bysubject_perElec.pkl', 'wb') as dataSet5sec:
            pickle.dump([df_train, df_test, df], dataSet5sec)

        '#print for validation:'
        print('all:', df.shape, 'N_On_lables:', df['labels'].sum())
        print('train:', df_train.shape, 'N_On_lables:', df_train['labels'].sum())
        print('test:', df_test.shape, 'N_On_lables:', df_test['labels'].sum())

        return df_train ,df_test,df


    else:

        with open('dataSet_zscore_bysubject_perElec.pkl','rb') as dataSet5sec:
            df_train, df_test, df = pickle.load(dataSet5sec)

        '#print for validation:'
        print('all:', df.shape, 'N_On_lables:', df['labels'].sum())
        print('train:', df_train.shape, 'N_On_lables:', df_train['labels'].sum())
        print('test:', df_test.shape, 'N_On_lables:', df_test['labels'].sum())


        return df_train ,df_test,df


def Seprates_Features_Labels(df_train,metaDataCol_N):


    labelData_2col=False

    df_xTrain = df_train.iloc[:, metaDataCol_N:]
    x_train = df_xTrain.values
    df_yTrain = df_train[["labels"]]
    if labelData_2col:
        df_yTrain.loc[:,'labels2'] = (1 - df_yTrain.loc[:,'labels'].copy())
    y_train = df_yTrain.values

    return x_train ,y_train,df_xTrain


def plot_Loss_accuracy(history):
    import matplotlib.pyplot as plt


    print(history.history.keys())

    fig = plt.figure()
    plt.plot(history.history['acc'])
    plt.plot(history.history['val_acc'])
    plt.title('model accuracy')
    plt.ylabel('accuracy')
    plt.xlabel('epoch')
    plt.legend(['train', 'val'], loc='upper left')
    fig.savefig('accuracy.png')

    # summarize history for loss
    fig = plt.figure()
    plt.plot(history.history['loss'])
    plt.plot(history.history['val_loss'])
    plt.title('model loss')
    plt.ylabel('loss')
    plt.xlabel('epoch')
    plt.legend(['train', 'val'], loc='upper left')
    # plt.show()
    fig.savefig('loss.png')

    # ____end visulize

    return

def Plot2PPT(s):
    from pptx import Presentation
    from pptx.util import Inches
    import os.path

    if os.path.exists(s.PPTName):
        prs = Presentation(s.PPTName)
    else:
        prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text ='Test: '+ str("%.2f" % s.score)+'(score),'+str("%.2f" % s.acc)+'(accuracy)'

    tf = body_shape.text_frame
    tf.text = 'model: '+s.title

    img_path1 = 'accuracy.png'
    img_path2 = 'loss.png'
    img_path3='model_plot.png'


    left = Inches(0)
    top = Inches(3)
    height = Inches(4)
    pic = slide.shapes.add_picture(img_path1, left, top, height=height)


    left = Inches(5)
    pic = slide.shapes.add_picture(img_path2, left, top, height=height)

    top=Inches(0.5)
    height= Inches(2.8)
    left = Inches(8)

    pic = slide.shapes.add_picture(img_path3, left, top, height=height)



    prs.save(s.PPTName)

    return

def aucPersubject(df,metaDataCol_N):
    from sklearn import metrics
    import pandas as pd

    df_groups = df.groupby([df['iVisitType']])

    iName=0
    auc_df=pd.DataFrame()
    for name, group in df_groups:
        iName=iName+1
        auc_df.loc[iName, 'subject'] = str(name[0]) + '_' + name[1]
        for col in group.iloc[:,metaDataCol_N:]:
            pred=group[col]
            y=group['labels']
            #remove index with NaNsS
            if pred.isnull().values.any():
                NanIndex = pred.apply(np.isnan)
                y=y.drop(pred.index[NanIndex])
                pred=pred.drop(pred.index[NanIndex])

            fpr, tpr, thresholds = metrics.roc_curve(y, pred, pos_label=1)
            aucResult=metrics.auc(fpr, tpr)
            if aucResult<0.5:
                aucResult=1-aucResult
            auc_df.loc[iName,col]=round(aucResult, 2)
            print('auc'+str(aucResult))
            del fpr, tpr, thresholds,pred,y,aucResult

    auc_df.loc['avg_all', :] = auc_df.mean(axis=0)
    auc_df.loc['std_all', :] = auc_df.std(axis=0)
    auc_df_t=auc_df.transpose()
    return auc_df_t

def combineRows2moreColumns(df):

    groups=df.groupby([df['Elec']])
    firstGroup_flag=1
    for name,group in groups:
        if firstGroup_flag==1:
            df_merge=group.copy()
            firstGroup_flag=0
            oldName=name
            print(name)
            print(df_merge.shape)

        else:
            if 'TrainOrTest' in df:
                df_merge=df_merge.merge( group, on=['reference','side','iVisit','OnOff','iPart','TrainOrTest','labels','iVisitType'],suffixes=('_'+oldName, '_'+name))#
            else:
                df_merge=df_merge.merge( group, on=['reference','side','iVisit','OnOff','iPart','labels','iVisitType'],suffixes=('_'+oldName, '_'+name))#
            oldName=name
            print(name)
            print(df_merge.shape)
    df_merge = df_merge.drop(['Elec_0_1','Elec_0_2','Elec_0_3','Elec_2_1','Elec_3_2','Elec_3_1',], 1)
    return df_merge


def add2Excel(C_paramter,trainAccuracy,testAccuracy,featureNum,modelName,line,df_results_new,subjectName):
    print(modelName+' -  Train:' + str(round(trainAccuracy, 2)) + '  Test:' + str(round(testAccuracy, 2)))
    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'C'] =(', '.join(str(x) for x in [C_paramter]))
    df_results_new.loc[line, 'featureSelection_N'] = featureNum
    df_results_new.loc[line, 'Model'] = modelName
    df_results_new.loc[line, 'AccuracyTrain'] = str(round(trainAccuracy, 2))
    df_results_new.loc[line, 'AccuracyTest'] = str(round(testAccuracy, 2))
    line = line + 1
    return line, df_results_new

def add2Excel_v2(C_paramter,Gama_paramter ,trainAccuracy, train_sensitivity, train_specificity,valAccuracy,val_sensitivity, val_specificity, testAccuracy,test_sensitivity, test_specificity, featureNum, modelName, line, df_results_new, subjectName, iter):


    print(modelName+' -  Train:' + str(round(trainAccuracy, 2)) + '  Test:' + str(round(testAccuracy, 2)))
    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    df_results_new.loc[line, 'C'] =(', '.join(str(x) for x in [C_paramter]))
    df_results_new.loc[line, 'Gamma'] =(', '.join(str(x) for x in [Gama_paramter]))
    df_results_new.loc[line, 'featureSelection_N'] = featureNum
    df_results_new.loc[line, 'Model'] = modelName
    df_results_new.loc[line, 'AccuracyTrain'] = str(round(trainAccuracy, 2))
    df_results_new.loc[line, 'train_sensitivity'] = str(round(train_sensitivity, 2))
    df_results_new.loc[line, 'train_specificity'] = str(round(train_specificity, 2))

    df_results_new.loc[line, 'AccuracyVal'] = str(round(valAccuracy, 2))
    df_results_new.loc[line, 'val_sensitivity'] = str(round(val_sensitivity, 2))
    df_results_new.loc[line, 'val_specificity'] = str(round(val_specificity, 2))

    df_results_new.loc[line, 'AccuracyTest'] = str(round(testAccuracy, 2))
    df_results_new.loc[line, 'test_sensitivity'] = str(round(test_sensitivity, 2))
    df_results_new.loc[line, 'test_specificity'] = str(round(test_specificity, 2))
    line = line + 1
    return line, df_results_new

def add2Excel_v2_gridSearch(C_paramter,Gama_paramter ,trainAccuracy,trainAUC, train_Precision, train_recall,testCVAccuracy,testCVAUC,testCV_Precision, testCV_recall, cleanValAccuracy,cleanVal_auc,cleanVal_Precision, cleanVal_recall, featureNum, modelName, line, df_results_new, subjectName,subject_number, iter):


    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectNumber'] = subject_number
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    df_results_new.loc[line, 'C'] =(', '.join(str(x) for x in [C_paramter]))
    df_results_new.loc[line, 'Gamma'] =(', '.join(str(x) for x in [Gama_paramter]))
    df_results_new.loc[line, 'featureSelection_N'] = str(featureNum)
    df_results_new.loc[line, 'Model'] = str(modelName)
    df_results_new.loc[line, 'AccuracyTrain'] = trainAccuracy
    df_results_new.loc[line, 'train_AUC'] = trainAUC
    df_results_new.loc[line, 'train_Precision'] =train_Precision
    df_results_new.loc[line, 'train_Recall'] = train_recall

    df_results_new.loc[line, 'testCV_accuracy'] = testCVAccuracy
    df_results_new.loc[line, 'testCV_AUC'] = testCVAUC
    df_results_new.loc[line, 'testCV_Precision'] = testCV_Precision
    df_results_new.loc[line, 'testCV_Recall'] = testCV_recall

    df_results_new.loc[line, 'cleanVal_accuracy'] = cleanValAccuracy
    df_results_new.loc[line, 'cleanVal_AUC'] = cleanVal_auc

    df_results_new.loc[line, 'cleanVal_Precision'] = cleanVal_Precision
    df_results_new.loc[line, 'cleanVal_Recall'] = cleanVal_recall
    line = line + 1
    return line, df_results_new




def add2Excel_RocLOO (results, featureNum, modelName, line, df_results_new, subjectName, subject_number, iter):
    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectNumber'] = subject_number
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    df_results_new.loc[line, 'featureSelection_N'] = str(featureNum)
    df_results_new.loc[line, 'Model'] = str(modelName)

    #train
    df_results_new.loc[line, 'train_accuracy_mean'] = results.train_accuracy.mean()
    df_results_new.loc[line, 'train_accuracy_std'] = results.train_accuracy.std()
    df_results_new.loc[line, 'train_precision_mean'] = results.train_precision.mean()
    df_results_new.loc[line, 'train_precision_std'] = results.train_precision.std()
    df_results_new.loc[line, 'train_recall_mean'] = results.train_recall.mean()
    df_results_new.loc[line, 'train_recall_std'] = results.train_recall.std()
    df_results_new.loc[line, 'train_roc_auc_mean'] = results.train_roc_auc.mean()
    df_results_new.loc[line, 'train_roc_auc_std'] = results.train_roc_auc.std()
    df_results_new.loc[line, 'train_senstivity_mean'] = results.train_senstivity.mean()
    df_results_new.loc[line, 'train_senstivity_std'] = results.train_senstivity.std()
    df_results_new.loc[line, 'train_specificity_mean'] = results.train_specificity.mean()
    df_results_new.loc[line, 'train_specificity_std'] = results.train_specificity.std()
    #test
    df_results_new.loc[line, 'test_accuracy'] = results.test_accuracy
    df_results_new.loc[line, 'test_recall'] = results.test_recall
    df_results_new.loc[line, 'test_precision'] = results.test_precision


    return line, df_results_new

def add2Excel_RocCV (results, featureNum, modelName, line, df_results_new, subjectName, subject_number, iter):
    clean_val_flag=1
    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectNumber'] = subject_number
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    df_results_new.loc[line, 'featureSelection_N'] = str(featureNum)
    df_results_new.loc[line, 'Model'] = str(modelName)

    #train
    df_results_new.loc[line, 'train_accuracy_mean'] = results.train_accuracy.mean()
    df_results_new.loc[line, 'train_accuracy_std'] = results.train_accuracy.std()
    df_results_new.loc[line, 'train_precision_mean'] = results.train_precision.mean()
    df_results_new.loc[line, 'train_precision_std'] = results.train_precision.std()
    df_results_new.loc[line, 'train_recall_mean'] = results.train_recall.mean()
    df_results_new.loc[line, 'train_recall_std'] = results.train_recall.std()
    df_results_new.loc[line, 'train_roc_auc_mean'] = results.train_roc_auc.mean()
    df_results_new.loc[line, 'train_roc_auc_std'] = results.train_roc_auc.std()
    df_results_new.loc[line, 'train_senstivity_mean'] = results.train_senstivity.mean()
    df_results_new.loc[line, 'train_senstivity_std'] = results.train_senstivity.std()
    df_results_new.loc[line, 'train_specificity_mean'] = results.train_specificity.mean()
    df_results_new.loc[line, 'train_specificity_std'] = results.train_specificity.std()
    #test
    df_results_new.loc[line, 'test_accuracy_mean'] = results.test_accuracy.mean()
    df_results_new.loc[line, 'test_accuracy_std'] = results.test_accuracy.std()
    df_results_new.loc[line, 'test_recall_mean'] = results.test_recall.mean()
    df_results_new.loc[line, 'test_recall_std'] = results.test_recall.std()
    df_results_new.loc[line, 'test_precision_mean'] = results.test_precision.mean()
    df_results_new.loc[line, 'test_precision_std'] = results.test_precision.std()
    df_results_new.loc[line, 'test_AUC_mean'] = results.test_auc.mean()
    df_results_new.loc[line, 'test_AUC_std'] = results.test_auc.std()

    if clean_val_flag:
        #cleanVal
        df_results_new.loc[line, 'cleanVal_accuracy'] = results.cleanVal_accuracy
        df_results_new.loc[line, 'cleanVal_recall'] = results.cleanVal_recall
        df_results_new.loc[line, 'cleanVal_precision'] = results.cleanVal_precision

    return line, df_results_new



def add2Excel_v3_gridSearch (C_paramter, Gama_paramter, results,scorsList, featureNum, modelName, line, df_results_new, subjectName, subject_number, iter):

    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectNumber'] = subject_number
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    df_results_new.loc[line, 'C'] =(', '.join(str(x) for x in [C_paramter]))
    df_results_new.loc[line, 'Gamma'] =(', '.join(str(x) for x in [Gama_paramter]))
    df_results_new.loc[line, 'featureSelection_N'] = str(featureNum)
    df_results_new.loc[line, 'Model'] = str(modelName)

    for  score in scorsList:

        trainName = 'train_' + score
        df_results_new.loc[line, trainName+'_mean'] = results[trainName].mean()
        df_results_new.loc[line, trainName+'_std'] = results[trainName].std()

        testName = 'test_' + score
        df_results_new.loc[line, testName+'_mean'] = results[testName].mean()
        df_results_new.loc[line, testName+'_std'] = results[testName].std()

    line = line + 1
    return line, df_results_new

def add2Excel_v4_gridSearch(iRandNum,psd_num,burst_num ,trainAccuracy,trainAUC, train_Precision, train_recall,trainCohen,trainSpecificity,testCVAccuracy,testCVAUC,testCV_Precision, testCV_recall,testCVCohen,testCVSpecificity, cleanValAccuracy,cleanVal_auc,cleanVal_Precision, cleanVal_recall,cleanVal_Cohen,cleanVal_specificity, featureNum, modelName, line, df_results_new, subjectName,subject_number, iter):


    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectNumber'] = subject_number
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    df_results_new.loc[line, 'iRandNum/trainNvisit'] = iRandNum
    df_results_new.loc[line, 'psd_num'] =(', '.join(str(x) for x in [psd_num]))
    df_results_new.loc[line, 'burst_num'] =(', '.join(str(x) for x in [burst_num]))
    df_results_new.loc[line, 'featureSelection_N'] = str(featureNum)
    df_results_new.loc[line, 'Model'] = str(modelName)

    df_results_new.loc[line, 'AccuracyTrain'] = trainAccuracy
    df_results_new.loc[line, 'train_AUC'] = trainAUC
    df_results_new.loc[line, 'train_Precision'] =train_Precision
    df_results_new.loc[line, 'train_Recall'] = train_recall
    df_results_new.loc[line, 'trainCohen'] = trainCohen
    df_results_new.loc[line, 'trainSpecificity'] = trainSpecificity

    df_results_new.loc[line, 'testCV_accuracy'] = testCVAccuracy
    df_results_new.loc[line, 'testCV_AUC'] = testCVAUC
    df_results_new.loc[line, 'testCV_Precision'] = testCV_Precision
    df_results_new.loc[line, 'testCV_Recall'] = testCV_recall
    df_results_new.loc[line, 'testCVCohen'] = testCVCohen
    df_results_new.loc[line, 'testCVSpecificity'] = testCVSpecificity

    df_results_new.loc[line, 'cleanVal_accuracy'] = cleanValAccuracy
    df_results_new.loc[line, 'cleanVal_AUC'] = cleanVal_auc

    df_results_new.loc[line, 'cleanVal_Precision'] = cleanVal_Precision
    df_results_new.loc[line, 'cleanVal_Recall'] = cleanVal_recall
    df_results_new.loc[line, 'cleanVal_Cohen'] = cleanVal_Cohen
    df_results_new.loc[line, 'cleanVal_specificity'] = cleanVal_specificity
    line = line + 1
    return line, df_results_new

def add2Excel_v2_DL(modelH_Params, line,  df_results_new, subjectName, iter):

    strName = (', '.join(str(x) for x in [subjectName]))
    print(strName)
    df_results_new.loc[line, 'subjectName'] = strName
    df_results_new.loc[line, 'iter'] = iter
    if hasattr(modelH_Params, 'Modelname'):
        df_results_new.loc[line, 'Model'] = modelH_Params.Modelname
    if hasattr(modelH_Params, 'featureSelection_N'):
        df_results_new.loc[line, 'featureSelection_N'] = modelH_Params.featureSelection_N
    if hasattr(modelH_Params, 'optimzerName'):
        df_results_new.loc[line, 'optimzerName'] = modelH_Params.optimzerName
    if hasattr(modelH_Params, 'epochs'):
        df_results_new.loc[line, 'epochs'] = modelH_Params.epochs
    if hasattr(modelH_Params, 'Reg'):
        df_results_new.loc[line, 'Reg'] = modelH_Params.Reg
    if hasattr(modelH_Params, 'batch_size_n'):
        df_results_new.loc[line, 'batch_size_n'] = modelH_Params.batch_size_n
    if hasattr(modelH_Params, 'activationName'):
        df_results_new.loc[line, 'activationName'] = modelH_Params.activationName
    if hasattr(modelH_Params, 'lossfunction'):
        df_results_new.loc[line, 'lossfunction'] = modelH_Params.lossfunction
    if hasattr(modelH_Params, 'dropoutV'):
        df_results_new.loc[line, 'dropoutV'] = modelH_Params.dropoutV

    df_results_new.loc[line, 'Train score'] = str(modelH_Params.Train_score)
    df_results_new.loc[line, 'Train accuracy'] = str(modelH_Params.Train_accuracy)

    df_results_new.loc[line, 'Val score'] = str(modelH_Params.Val_score)
    df_results_new.loc[line, 'Val accuracy'] = str(modelH_Params.Val_accuracy)

    df_results_new.loc[line, 'Test score'] = str(modelH_Params.Test_score)
    df_results_new.loc[line, 'Test accuracy'] = str(modelH_Params.Test_accuracy)

    line = line + 1

    return line, df_results_new




def plot_coefficients(classifier, feature_names,resultPath,iteration,top_features=5):
    import matplotlib.pyplot as plt
    import os


    coef =classifier.best_estimator_.coef_.ravel()
    top_positive_coefficients = np.argsort(coef)[-top_features:]
    top_negative_coefficients = np.argsort(coef)[:top_features]
    top_coefficients = np.hstack([top_negative_coefficients, top_positive_coefficients])
    print('plot coefficeients')

    plt.figure
    colors = ["red" if c < 0 else "blue" for c in coef[top_coefficients]]
    plt.bar(np.arange(2 * top_features), coef[top_coefficients], color=colors)
    feature_names = np.array(feature_names)
    plt.xticks(np.arange(0, 2 * top_features), feature_names[top_coefficients],fontsize=6, rotation=40, ha="right")
    plt.show()
    Path=resultPath+'/coefficientsFigures/'
    if not os.path.exists(Path):
        os.makedirs(Path)

    plt.savefig(Path+"coefficients_"+str(iteration)+".pdf")
    return



def plot_coefficients_v2(classifier, feature_names,resultPath,iteration,figTitle,top_features=2):
    from matplotlib.font_manager import FontProperties
    import matplotlib.pyplot as plt
    import os


    coef =classifier.best_estimator_.coef_.ravel()
    top_positive_coefficients = np.argsort(coef)[-top_features:]
    top_negative_coefficients = np.argsort(coef)[:top_features]
    top_coefficients = np.hstack([top_negative_coefficients, top_positive_coefficients])
    print('plot coefficeients')

    plt.figure(figsize=(30, 30))
    plt.figure
    colors = ["red" if c < 0 else "blue" for c in coef[top_coefficients]]
    plt.bar(np.arange(2 * top_features), coef[top_coefficients], color=colors)
    feature_names = np.array(feature_names)
    plt.xticks(np.arange(0, 2 * top_features), feature_names[top_coefficients], fontsize=18, rotation=20, ha="right")
    # plt.show()
    plt.title(figTitle)
    Path = resultPath + '/coefficientsFigures/'
    if not os.path.exists(Path):
        os.makedirs(Path)
    plt.savefig(Path + "coefficients_" + str(iteration) + ".pdf")
    return


def plot_abs_coefficients_v2(classifier, feature_names,resultPath,iteration,figTitle,top_features=2):
    from matplotlib.font_manager import FontProperties
    import matplotlib.pyplot as plt
    import os

    coef_abs =abs(classifier.best_estimator_.coef_.ravel())
    top_coefficients = np.argsort(coef_abs)[-top_features:]
    print('plot coefficeients')

    plt.figure(figsize=(45, 15))
    plt.figure.autolayout: True
    # colors = ["red" if c < 0 else "blue" for c in coef_abs[top_coefficients]]
    colors = ["blue"]
    plt.barh(np.arange(top_features), coef_abs[top_coefficients], color=colors)
    feature_names = np.array(feature_names)
    plt.yticks(np.arange(0, top_features), feature_names[top_coefficients], fontsize=95, ha="right")
    # plt.show()
    plt.xticks(fontsize=65, rotation=270)
    plt.xlim(0, 1.1)
    plt.title(figTitle, fontsize=50)
    plt.tick_params(axis='y', which='major', pad=50)

    plt.tight_layout()

    Path = resultPath + '/absCoefficientsFigures/'
    if not os.path.exists(Path):
        os.makedirs(Path)
    plt.savefig(Path + "coefficients_abs" + str(iteration) + ".pdf")
    # plt.savefig(Path + "coefficients_abs" + str(iteration) + ".PNG")

    return


def count_psd_vs_burst_number(fea_name):
    #Descriptuion this func is counting how many power psd features are they comperd to burst features
    psd_num = 0
    burst_num = 0
    for i in fea_name :
       # temp= i[-7:]
        if 'MeanPSD' in i:
            psd_num   = psd_num+1
        else:
            burst_num = burst_num+1


    return psd_num,burst_num